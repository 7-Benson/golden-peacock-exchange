#!/usr/bin/env python3
"""Generate Golden Peacock Exchange PPT - Multi-Asset Compliance Platform"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
GOLD = RGBColor(0xD4, 0xA5, 0x37)
GOLD_LIGHT = RGBColor(0xF0, 0xD0, 0x60)
GOLD_DARK = RGBColor(0xB8, 0x86, 0x0B)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
DARK2 = RGBColor(0x16, 0x21, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_BLUE = RGBColor(0x54, 0x9B, 0xFF)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)


def set_slide_bg(slide, color=DARK2):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_gold_accent_line(slide, left, top, width, height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body, title_color=GOLD, bg_color=RGBColor(0x22, 0x2E, 0x44)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = GOLD_DARK
    card.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(0.9),
                 body, font_size=13, color=LIGHT_GRAY)
    return card


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_char="▸"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(8)
    return tb


def add_section_number(slide, num, total=16):
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=GOLD_DARK, alignment=PP_ALIGN.RIGHT)


def add_header_bar(slide, title, subtitle=None, page_num=None):
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), GOLD)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 title, font_size=32, color=GOLD, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=LIGHT_GRAY)
    add_gold_accent_line(slide, Inches(0.8), Inches(1.4), Inches(3), Pt(2))
    if page_num:
        add_section_number(slide, page_num)


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Microsoft YaHei'
                paragraph.font.color.rgb = WHITE if r == 0 else LIGHT_GRAY
                if r == 0:
                    paragraph.font.bold = True
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GOLD_DARK
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3E) if r % 2 == 1 else RGBColor(0x22, 0x2E, 0x44)
    return table_shape


TOTAL_SLIDES = 16

# ═══════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(2.8), RGBColor(0x12, 0x1A, 0x30))
add_shape_bg(slide, Inches(0), Inches(2.8), Inches(13.333), Pt(5), GOLD)

# Peacock emblem
circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.0), Inches(2.3), Inches(2.3))
circ.fill.solid()
circ.fill.fore_color.rgb = GOLD
circ.line.fill.background()
circ2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.4), Inches(1.5), Inches(1.5))
circ2.fill.solid()
circ2.fill.fore_color.rgb = GOLD_DARK
circ2.line.fill.background()
circ3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.4), Inches(1.9), Inches(0.5), Inches(0.5))
circ3.fill.solid()
circ3.fill.fore_color.rgb = WHITE
circ3.line.fill.background()

add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.2),
             "金孔雀交易所", font_size=54, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8),
             "GOLDEN PEACOCK EXCHANGE", font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5), Inches(5.3), Inches(3.333), Pt(3), GOLD)

add_text_box(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.6),
             "综合性合规资产交易平台  —  GPC平台通证首发", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
             "链接多元价值 · 开启资产新纪元", font_size=16, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(7.1), Inches(13.333), Pt(4), GOLD)
add_section_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 2: PROJECT OVERVIEW
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "项目概览", "Project Overview", 2)

# Positioning
add_shape_bg(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8), RGBColor(0x22, 0x2E, 0x44))
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.6),
             "多种合规复合型资产的综合性交易平台 — 基于币安链（BNB Chain）构建，原生平台通证GPC驱动生态",
             font_size=18, color=GOLD_LIGHT, bold=True)

# Three key cards
cards_data = [
    ("🏛️", "平台定位", "综合性合规资产交易平台\n支持股票证券/数字资产\n质押/Launchpad/RWA"),
    ("💎", "平台通证 GPC", "总发行量 1,000,000,000\nBEP-20标准\n驱动平台生态价值"),
    ("🔮", "战略布局", "多元合规资产交易\n股票证券+RWA+\n数字资产三位一体"),
]
for i, (icon, title, body) in enumerate(cards_data):
    left = Inches(0.8 + i * 4.1)
    add_card(slide, left, Inches(2.8), Inches(3.7), Inches(2.0), f"{icon} {title}", body)

# Vision
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
             "项目愿景", font_size=20, color=GOLD, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.6),
             "构建安全、合规、透明的多元化资产交易生态，让全球用户自由参与多种合规资产的价值流通",
             font_size=17, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 3: MARKET OPPORTUNITY
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "市场机遇", "Market Opportunity", 3)

cards3 = [
    ("📊", "合规资产数字化浪潮",
     "• 传统股票证券数字化交易规模持续增长\n• 机构资金加速入场\n• 全球合规化趋势明显\n• 多元化资产配置需求旺盛"),
    ("🌏", "东南亚多元金融浪潮",
     "• 东南亚加密采用率全球领先\n• 年轻人口+高移动渗透率\n• 多元化合规资产需求爆发\n→ 数字金融弯道超车"),
    ("🏆", "平台竞争壁垒",
     "• 合规+多元化资产+本地化 = 核心壁垒\n• 单一资产交易平台已不满足需求\n• GPC平台币驱动生态飞轮\n• 股票证券+数字资产+RWA三位一体"),
]
for i, (icon, title, body) in enumerate(cards3):
    left = Inches(0.8 + i * 4.1)
    add_card(slide, left, Inches(1.8), Inches(3.7), Inches(2.8), f"{icon} {title}", body)

# Future note
add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.0), RGBColor(0x22, 0x2E, 0x44))
add_text_box(slide, Inches(1.0), Inches(5.05), Inches(11.3), Inches(0.3),
             "🌟 战略前瞻", font_size=16, color=GOLD, bold=True)
add_text_box(slide, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.5),
             "RWA（真实世界资产代币化）市场预计2030年达数万亿级别 — 金孔雀交易所将此作为核心战略方向之一",
             font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 4: PLATFORM VALUE PROPOSITION
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "平台核心价值", "Platform Core Value Proposition", 4)

props = [
    ("🛡️", "安全可靠", "多层安全防护体系\n冷热钱包分离\n第三方安全审计"),
    ("⚡", "高效流畅", "基于BNB Chain\n高吞吐低延迟\n毫秒级交易确认"),
    ("🔍", "透明可信", "链上资产可查\n定期审计报告\n合规KYC/AML"),
    ("🌐", "多元合规", "覆盖股票证券\n数字资产、RWA\n多资产类别合规交易"),
]
for i, (icon, title, desc) in enumerate(props):
    left = Inches(0.8 + i * 3.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.8), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x22, 0x2E, 0x44)
    card.line.color.rgb = GOLD_DARK
    card.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.2), Inches(1.9), Inches(2.4), Inches(0.5),
                 f"{icon} {title}", font_size=20, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(2.5), Inches(2.4), Inches(1.5),
                 desc, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Value chain
add_card(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.5),
         "平台生态飞轮",
         "GPC平台币价值 → 赋能交易/质押/治理 → 吸引用户与流动性 → 生态繁荣驱动GPC需求 → GPC价值提升 → 反哺平台发展")


# ═══════════════════════════════════════════════
# SLIDE 5: PLATFORM CORE FEATURES
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "平台核心功能", "Core Platform Features", 5)

feats = [
    ("💹", "股票/证券交易", "合规股票证券交易通道\n多元化投资选择\n主流市场接入"),
    ("⛏️", "GPC质押挖矿", "质押GPC获取收益\n灵活质押期限\n按月手续费结算"),
    ("🚀", "Launchpad", "优质项目首发平台\nGPC持有者优先认购\n代币销售与拍卖"),
    ("🔗", "RWA资产发行", "现实资产代币化\n合规发行交易\n多类型资产支持"),
]
for i, (icon, title, desc) in enumerate(feats):
    left = Inches(0.8 + i * 3.1)
    add_card(slide, left, Inches(1.8), Inches(2.8), Inches(2.3), f"{icon} {title}", desc)

# More features
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.4),
             "更多功能", font_size=18, color=GOLD, bold=True)
more_feats = [
    "数字资产现货交易 — 主流交易对，订单簿深度，专业级交易界面",
    "钱包服务 — 安全资产管理，多链支持",
    "行情数据 — 实时K线、深度图、市场分析",
    "API接口 — 面向量化交易与机构客户",
]
add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.5), more_feats, font_size=14)


# ═══════════════════════════════════════════════
# SLIDE 6: TECHNOLOGY ARCHITECTURE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "技术架构", "Technology Architecture", 6)

# Blockchain
add_card(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.5),
         "底层公链 — 币安链（BNB Chain）",
         "BEP-20代币标准 | 高吞吐量（~300 TPS）| 低交易费用（<0.1 USD）\n"
         "成熟的DeFi生态（PancakeSwap等DEX协议） | 全球40+验证节点 | 开发者生态系统完善")

# Core modules
modules = [
    ("💱", "交易引擎", "高性能撮合引擎\n订单簿深度管理\n实时行情推送"),
    ("⛓️", "智能合约", "质押挖矿合约\nLaunchpad合约\nRWA发行标准"),
    ("🛡️", "安全系统", "KYC/AML认证系统\n风控引擎\n冷热钱包分离"),
    ("💳", "资产管理系统", "多链钱包支持\n证券资产托管\n充值提现系统"),
]
for i, (icon, title, desc) in enumerate(modules):
    left = Inches(0.8 + i * 3.1)
    add_card(slide, left, Inches(3.6), Inches(2.8), Inches(2.0), f"{icon} {title}", desc)

# Security
add_card(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.7),
         "安全审计", "第三方智能合约代码审计  |  币安链生态安全标准  |  多重签名钱包  |  24/7安全监控")


# ═══════════════════════════════════════════════
# SLIDE 7: GPC TOKENOMICS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "GPC代币经济模型", "GPC Tokenomics", 7)

# Token info
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(1.0),
         "代币信息",
         "代币符号：GPC  |  总发行量：1,000,000,000  |  标准：BEP-20  |  公链：BNB Chain")

# Allocation table
table_data = [
    ["用途", "比例", "数量", "锁仓规则"],
    ["🌱 生态与社区激励", "60%", "6亿", "按智能合约逐步释放"],
    ["🏦 生态基金/战略储备", "15%", "1.5亿", "基金会多签管理"],
    ["💧 流动性储备", "10%", "1亿", "上线即部分流通"],
    ["👥 团队与创始人", "2%", "2000万", "锁36个月，分36月线性释放"],
    ["🌟 创世节点", "1%", "1000万", "锁36个月"],
    ["📢 市场推广", "5%", "5000万", "按活动计划释放"],
    ["🏛️ 基金会储备", "7%", "7000万", "多签钱包管理"],
]

add_table(slide, Inches(0.8), Inches(3.1), Inches(7.5), Inches(3.2),
          len(table_data), 4, table_data,
          col_widths=[Inches(2.3), Inches(1.2), Inches(1.2), Inches(2.8)])

# Summary callout
add_shape_bg(slide, Inches(9.0), Inches(1.8), Inches(3.8), Inches(4.5), RGBColor(0x1A, 0x2A, 0x3E))
add_text_box(slide, Inches(9.2), Inches(1.9), Inches(3.4), Inches(0.4),
             "📊 关键数据", font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
key_data = [
    "总发行量",
    "1,000,000,000",
    "GPC",
    "",
    "公链",
    "BNB Chain (BEP-20)",
    "",
    "生态占比",
    "75%",
    "(生态60%+储备15%)",
]
for i, line in enumerate(key_data):
    if not line:
        continue
    fs = 20 if i in [1, 8] else 14
    clr = GOLD if i in [1, 8] else (GOLD_LIGHT if i in [0, 4, 7] else LIGHT_GRAY)
    bld = True if i in [0, 1, 4, 7, 8] else False
    add_text_box(slide, Inches(9.2), Inches(2.5 + i * 0.35), Inches(3.4), Inches(0.35),
                 line, font_size=fs, color=clr, bold=bld, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 8: GPC TOKEN UTILITY
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "GPC代币价值与用途", "GPC Token Utility & Value Capture", 8)

utilities = [
    ("💱", "手续费折扣", "持有GPC可享受\n平台交易手续费减免\n持有越多折扣越大"),
    ("💰", "合伙人分润", "质押GPC成为平台合伙人\n按月分享交易所手续费\n持有越多分润越高"),
    ("🗳️", "平台治理", "GPC持有者参与\n平台重大决策投票\n社区自治管理"),
    ("🚀", "Launchpad认购", "优质项目的优先\n认购权与额度分配\nGPC持有者专属"),
    ("🏪", "生态支付", "GPC作为平台内\n流通与支付介质\n交易对基础货币"),
    ("📈", "价值增值", "多元化资产生态增长驱动\nGPC需求持续上升\n通缩机制持续赋能"),
]
for i, (icon, title, desc) in enumerate(utilities):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.8 + row * 2.5)
    add_card(slide, left, top, Inches(3.7), Inches(2.0), f"{icon} {title}", desc)


# ═══════════════════════════════════════════════
# SLIDE 9: TOKEN DISTRIBUTION SCHEDULE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "代币释放计划", "Token Release Schedule", 9)

# Schedule cards
schedules = [
    ("👥 团队与创始人", "2%", "锁仓36个月", "第1-36个月线性释放\n每月释放约0.056%"),
    ("🌟 创世节点", "1%", "锁仓36个月", "第1-36个月分批释放\n持有人锁仓周期36个月"),
    ("🌱 生态与社区激励", "60%", "按智能合约逐步释放", "交易挖矿/质押奖励/\n社区空投分期释放"),
    ("其他", "37%", "分批释放", "流动性/基金/市场/\n储备按计划管理"),
]
for i, (title, pct, lock, release) in enumerate(schedules):
    left = Inches(0.8 + i * 3.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.8), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3E)
    card.line.color.rgb = GOLD
    card.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.1), Inches(1.9), Inches(2.6), Inches(0.4),
                 title, font_size=15, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.3), Inches(2.6), Inches(0.4),
                 pct, font_size=22, color=GOLD_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.7), Inches(2.6), Inches(0.3),
                 f"🔒 {lock}", font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(3.1), Inches(2.6), Inches(0.8),
                 release, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Timeline visual
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8), RGBColor(0x22, 0x2E, 0x44))
add_text_box(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.3),
             "释放时间线（示意）", font_size=16, color=GOLD, bold=True)

timeline_items = [
    ("T+0", "流动性储备\n市场推广部分", "上线即部分流通"),
    ("T+持续", "生态激励\n逐步释放", "按智能合约规则"),
    ("T+36月", "创世节点/团队\n开始分批释放", "36个月逐步解锁"),
]
for i, (time, who, detail) in enumerate(timeline_items):
    left = Inches(0.8 + i * 3.1 + 0.7)
    add_text_box(slide, left, Inches(5.3), Inches(2.8), Inches(0.3),
                 f"📌 {time}", font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(5.6), Inches(2.8), Inches(0.3),
                 who, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(5.9), Inches(2.8), Inches(0.4),
                 detail, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 10: PARTNER REVENUE SHARING
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "合伙人收益计划", "Partner Revenue Sharing", 10)

# Core concept
add_card(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2),
         "核心理念",
         "持有并质押GPC，即可成为金孔雀交易所合伙人，享有双重收益：全网提币手续费70%分红 + 5代下级交易手续费返佣。")

# Flow diagram - visual using cards
flow_cards = [
    ("💹", "交易所产生手续费", "提币手续费+交易手续费", Inches(0.8)),
    ("→", "双重收益分配", "70%提币费+5代返佣", Inches(4.3)),
    ("→", "月结汇总", "按质押比例+层级分配", Inches(6.8)),
    ("💰", "合伙人钱包", "每月自动到账", Inches(9.3)),
]
for i, (icon, title, desc, left) in enumerate(flow_cards):
    if i in [1, 2]:
        add_text_box(slide, left - Inches(0.5), Inches(3.3), Inches(1.0), Inches(0.6),
                     icon, font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
        add_card(slide, left, Inches(3.1), Inches(2.8), Inches(1.0), title, desc, bg_color=RGBColor(0x2A, 0x3A, 0x50))
    else:
        add_card(slide, left, Inches(3.1), Inches(3.2), Inches(1.0), f"{icon} {title}", desc)

# Two revenue streams
add_shape_bg(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(1.8), RGBColor(0x1A, 0x2A, 0x3E))
add_text_box(slide, Inches(1.0), Inches(4.6), Inches(5.1), Inches(0.3),
             "💰 双重收益体系", font_size=16, color=GOLD, bold=True)
streams = [
    "收益一：全网提币手续费的70%按GPC质押比例分配",
    "收益二：5代以内下级用户交易手续费返佣",
    "合伙人锁仓周期：36个月",
    "手续费结算方式：按月结算，自动到账",
]
add_bullet_list(slide, Inches(1.0), Inches(5.0), Inches(5.1), Inches(1.2), streams, font_size=13)

# Two tiers
tiers = [
    ("🌟 创世合伙人", "• 创世节点持有者（1%配额）\n• 最高分润比例\n• 享70%提币费+5代返佣\n• 锁仓36个月\n• 平台早期治理权", True),
    ("💎 GPC质押合伙人", "• 任何持有并质押GPC的用户\n• 按质押数量分配\n• 享70%提币费+5代返佣\n• 质押越多，分润越高\n• 锁仓36个月", False),
]
for i, (title, body, is_genesis) in enumerate(tiers):
    left = Inches(0.8 + i * 6.2)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.0), Inches(5.8), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3E)
    card.line.color.rgb = GOLD if is_genesis else GOLD_DARK
    card.line.width = Pt(2) if is_genesis else Pt(1)
    add_text_box(slide, left + Inches(0.3), Inches(5.1), Inches(5.2), Inches(0.4),
                 title, font_size=18, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.3), Inches(5.5), Inches(5.2), Inches(1.3),
                 body, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 11: MULTI-ASSET LAYOUT
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "多元资产战略布局", "Multi-Asset Strategy", 11)

# Strategic vision
add_card(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.0),
         "战略定位",
         "金孔雀交易所定位为综合性合规资产交易平台，实现三大资产类别的有机融合 — 一个平台，多种合规资产选择")

# Three asset cards
asset_cards = [
    ("💹", "股票证券资产", "合规通道接入\n主流市场股票/ETF/证券交易\n多元化投资选择\n遵循各国证券法规"),
    ("⚡", "数字资产", "主流加密货币现货交易\nBNB Chain生态集成\n安全快捷的数字资产管理\n合规KYC/AML体系"),
    ("🏠", "RWA现实资产", "房地产等实体资产代币化\n合规发行与交易平台\n多类型RWA资产支持\n东南亚区域优先布局"),
]
for i, (icon, title, desc) in enumerate(asset_cards):
    left = Inches(0.8 + i * 4.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.1), Inches(3.7), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3E)
    card.line.color.rgb = GOLD
    card.line.width = Pt(1)
    add_text_box(slide, left, Inches(3.2), Inches(3.7), Inches(0.5),
                 f"{icon} {title}", font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(3.8), Inches(3.3), Inches(1.5),
                 desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Disclaimer
add_shape_bg(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.6), RGBColor(0x30, 0x20, 0x20))
add_text_box(slide, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.4),
             "⚖ 所有资产类别均严格遵循合规框架，持合法牌照运营，用户资产安全第一",
             font_size=14, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 12: ROADMAP
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "发展路线图", "Roadmap", 12)

phases = [
    ("Phase 1", "Q2-Q3 2026", "平台搭建与上线",
     [
         "完成平台技术开发",
         "GPC代币合约部署",
         "获取合规资质",
         "内部测试网运行",
     ]),
    ("Phase 2", "Q4 2026", "正式运营启动",
     [
         "主网上线与交易平台开放",
         "GPC上线交易",
         "Launchpad功能上线",
         "全球社区建设启动",
     ]),
    ("Phase 3", "2027", "生态扩张",
     [
         "多区域拓展",
         "股票证券交易通道上线",
         "RWA资产试点启动",
         "用户规模快速增长",
     ]),
    ("Phase 4", "2028+", "多元资产全面布局",
     [
         "多类型RWA资产全面开放",
         "全面合规证券交易体系",
         "成为东南亚多元资产交易标杆",
         "建立行业合规标准",
     ]),
]
for i, (phase, time, title, items) in enumerate(phases):
    left = Inches(0.5 + i * 3.2)
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.9), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = GOLD_DARK if i < 2 else RGBColor(0x2A, 0x3A, 0x50)
    header.line.fill.background()
    add_text_box(slide, left, Inches(1.85), Inches(2.9), Inches(0.3),
                 f"{phase} | {time}", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.2), Inches(2.9), Inches(0.3),
                 title, font_size=12, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, left + Inches(0.2), Inches(2.9), Inches(2.6), Inches(2.5),
                    items, font_size=12, color=LIGHT_GRAY, bullet_char="•")
    if i < len(phases) - 1:
        add_text_box(slide, left + Inches(3.0), Inches(2.0), Inches(0.3), Inches(0.5),
                     "→", font_size=20, color=GOLD, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 13: COMPLIANCE & SECURITY
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "合规与安全", "Compliance & Security", 13)

# Compliance
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.0),
         "合规体系",
         "🛡️ 持牌运营目标\n🔍 严格KYC/AML体系\n📊 证券交易合规牌照\n🌐 国际合规标准对标")

# Security
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.0),
         "安全保障",
         "🔒 冷热钱包分离\n🤖 智能合约多重审计\n🕵️ 24/7风控监控\n✅ 多重签名管理")

# Protection layers
layers = [
    ("第一层", "账户安全", "二次验证(2FA)\n反钓鱼保护\n登录监控"),
    ("第二层", "交易安全", "风控引擎\n异常交易检测\n提现审核"),
    ("第三层", "资产安全", "冷钱包离线存储\n多签钱包\n保险机制"),
    ("第四层", "系统安全", "代码审计\n渗透测试\nDDoS防护"),
]
for i, (layer, title, desc) in enumerate(layers):
    left = Inches(0.8 + i * 3.1)
    add_card(slide, left, Inches(4.2), Inches(2.8), Inches(1.8), f"{layer}: {title}", desc)


# ═══════════════════════════════════════════════
# SLIDE 14: RISK CONTROL
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "风控体系", "Risk Control Framework", 14)

risk_layers = [
    ("第一层", "内部风控", "内控制度健全\n操作权限管理\n定期内部审计", ACCENT_GREEN),
    ("第二层", "技术风控", "智能合约审计\n链上监控预警\n安全事件响应", ACCENT_BLUE),
    ("第三层", "合规风控", "KYC/AML审核\n交易监控\n举报与处置", GOLD),
    ("第四层", "外部风控", "第三方审计\n法律顾问\n保险保障", ACCENT_RED),
]
for i, (layer, title, desc, color) in enumerate(risk_layers):
    left = Inches(0.8 + i * 3.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.8), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3E)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    add_text_box(slide, left, Inches(1.9), Inches(2.8), Inches(0.4),
                 layer, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.3), Inches(2.8), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.8), Inches(2.6), Inches(1.0),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Risk disclosure
add_shape_bg(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.5), RGBColor(0x22, 0x2E, 0x44))
add_text_box(slide, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.3),
             "📋 信息披露与透明度", font_size=16, color=GOLD, bold=True)
transparency = [
    "链上资产数据实时可查 — 所有GPC链上交易公开透明",
    "定期第三方审计报告 — 平台资金与代币储备定期披露",
    "用户资产独立托管 — 平台运营资金与用户资产严格分离",
    "多元资产专属风控 — 股票/数字资产/RWA分别设立独立风控模型",
]
add_bullet_list(slide, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.0), transparency, font_size=14)


# ═══════════════════════════════════════════════
# SLIDE 15: CORE TEAM (TBD)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK2)
add_header_bar(slide, "核心团队", "Core Team", 15)

team = [
    ("👤", "CEO / 创始人", "待定"),
    ("👤", "CTO", "待定"),
    ("👤", "首席合规官", "待定"),
    ("👤", "首席运营官", "待定"),
]
for i, (icon, title, desc) in enumerate(team):
    left = Inches(0.8 + i * 3.1)
    add_card(slide, left, Inches(1.8), Inches(2.8), Inches(2.0), f"{icon} {title}", desc)

# Advisory board
add_card(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(1.5),
         "顾问团队（拟邀）",
         "区块链 / 金融 / 法律领域知名专家与机构顾问\n具体名单待定")

add_card(slide, Inches(6.8), Inches(4.2), Inches(5.5), Inches(1.5),
         "技术合作伙伴（拟邀）",
         "币安链生态合作伙伴\n区块链安全审计机构\n做市商与流动性合作伙伴")

# Note
add_shape_bg(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5), RGBColor(0x22, 0x2E, 0x44))
add_text_box(slide, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.4),
             "团队成员信息持续更新中，敬请期待",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 16: CLOSING
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.0), RGBColor(0x12, 0x1A, 0x30))
add_shape_bg(slide, Inches(0), Inches(3.0), Inches(13.333), Pt(5), GOLD)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(0.5),
             "金孔雀交易所  |  GOLDEN PEACOCK EXCHANGE",
             font_size=22, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.6),
             "核心信息回顾", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

recap = [
    "✅ 基于BNB Chain的综合性合规资产交易平台",
    "✅ 原生平台通证 GPC（总发行量10亿枚，BEP-20标准）",
    "✅ 股票证券 · 数字资产 · RWA 三类合规资产一体化交易",
    "✅ 安全合规 · 透明可信 · 多元资产 · 全球普惠",
]
add_bullet_list(slide, Inches(2), Inches(2.1), Inches(9.3), Inches(1.0),
                recap, font_size=15, color=LIGHT_GRAY)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
             "合作邀约", font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.4), Inches(11.3), Inches(0.8),
             "欢迎机构投资者 · 做市商 · 技术伙伴 · 社区贡献者 洽谈合作",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(3), Inches(5.3), Inches(7.333), Inches(1.0), RGBColor(0x22, 0x2E, 0x44))
add_text_box(slide, Inches(3.5), Inches(5.4), Inches(6.333), Inches(0.3),
             "联系方式（待定）", font_size=15, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(5.75), Inches(6.333), Inches(0.4),
             "官网 · 邮箱 · Telegram · 微信商务",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
             "「链接多元价值 · 开启资产新纪元」",
             font_size=22, color=GOLD_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(7.1), Inches(13.333), Pt(4), GOLD)
add_section_number(slide, 16, TOTAL_SLIDES)


# ── SAVE ──
output_dir = "/Users/mac/.openclaw/workspace/projects/金孔雀交易所PPT"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "金孔雀交易所_Golden_Peacock_Exchange.pptx")
prs.save(output_path)
print(f"✅ PPT saved: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
